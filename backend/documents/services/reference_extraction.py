from __future__ import annotations

from collections import Counter
from dataclasses import asdict, dataclass
from io import BytesIO
import math
import re
from typing import Protocol

from docx import Document as DocxDocument
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation  # type: ignore

from .token_metrics import estimate_token_count


@dataclass(frozen=True)
class NormalizedSection:
    id: str
    type: str
    title: str
    locator: dict[str, object]
    content: str
    estimated_tokens: int
    section_kind: str
    keywords: list[str]
    themes: list[str]
    signals: dict[str, float]


@dataclass(frozen=True)
class ReferenceUnderstandingResult:
    summary: str
    dominant_themes: list[str]
    domain_terms: list[str]
    document_patterns: list[str]
    signal_scores: dict[str, float]
    low_signal: bool


@dataclass(frozen=True)
class NormalizedDocument:
    document_summary: str
    sections: list[NormalizedSection]
    analysis: ReferenceUnderstandingResult


class ReferenceExtractionError(ValueError):
    pass


class ReferenceExtractor(Protocol):
    def extract(self, payload: bytes) -> NormalizedDocument: ...


SUPPORTED_REFERENCE_EXTENSIONS = {
    ".docx": "DOCX",
    ".pptx": "PPTX",
    ".xlsx": "XLSX",
    ".txt": "TXT",
    ".md": "MD",
    ".pdf": "PDF",
}

MAX_SECTION_CHARS = 1800
SUMMARY_MAX_CHARS = 600
STOPWORDS = {
    "und",
    "oder",
    "der",
    "die",
    "das",
    "den",
    "dem",
    "des",
    "ein",
    "eine",
    "einer",
    "einem",
    "für",
    "mit",
    "auf",
    "von",
    "ist",
    "im",
    "in",
    "zu",
    "am",
    "an",
    "bei",
    "nach",
    "als",
    "werden",
    "wird",
    "durch",
    "the",
    "and",
    "for",
    "with",
    "from",
    "this",
    "that",
    "into",
    "eine",
}
THEME_PATTERNS: dict[str, tuple[str, ...]] = {
    "risk_opportunity": ("risiko", "risiken", "chance", "chancen", "bewertung", "auswirkung", "prävent"),
    "instruction": ("anweisung", "instruction", "durchführung", "arbeitsablauf", "umsetzung", "durchzuführ"),
    "workflow": ("prozess", "ablauf", "workflow", "prozessablauf", "schnittstelle", "schritt"),
    "responsibility": ("verantwort", "zuständig", "rolle", "aufgabe", "befugnis", "owner"),
    "record_control": ("aufzeichnung", "dokument", "lenkung", "nachweis", "ablage", "aufbewahrung", "revision"),
    "policy": ("politik", "policy", "grundsatz", "zielsetzung", "leitlinie"),
    "compliance": ("compliance", "audit", "norm", "iso", "scc", "gesetz", "anforderung"),
}


def infer_reference_file_type(filename: str) -> str:
    lowered = filename.lower().rsplit(".", 1)
    ext = f".{lowered[-1]}" if len(lowered) == 2 else ""
    return SUPPORTED_REFERENCE_EXTENSIONS.get(ext, "OTHER")


def serialize_normalized_document(document: NormalizedDocument) -> dict[str, object]:
    return {
        "document_summary": document.document_summary,
        "sections": [asdict(item) for item in document.sections],
        "analysis": asdict(document.analysis),
    }


def deserialize_normalized_document(payload: dict[str, object]) -> NormalizedDocument:
    sections_raw = payload.get("sections") if isinstance(payload, dict) else []
    sections: list[NormalizedSection] = []
    if isinstance(sections_raw, list):
        for item in sections_raw:
            if not isinstance(item, dict):
                continue
            sections.append(
                NormalizedSection(
                    id=str(item.get("id", "")),
                    type=str(item.get("type", "section")),
                    title=str(item.get("title", "")),
                    locator=item.get("locator") if isinstance(item.get("locator"), dict) else {},
                    content=str(item.get("content", "")),
                    estimated_tokens=int(item.get("estimated_tokens", 0) or 0),
                    section_kind=str(item.get("section_kind", "section")),
                    keywords=[
                        str(keyword).strip()
                        for keyword in item.get("keywords", [])
                        if str(keyword).strip()
                    ]
                    if isinstance(item.get("keywords"), list)
                    else [],
                    themes=[
                        str(theme).strip()
                        for theme in item.get("themes", [])
                        if str(theme).strip()
                    ]
                    if isinstance(item.get("themes"), list)
                    else [],
                    signals={
                        str(key): float(value)
                        for key, value in item.get("signals", {}).items()
                    }
                    if isinstance(item.get("signals"), dict)
                    else {},
                )
            )
    analysis_raw = payload.get("analysis") if isinstance(payload, dict) else {}
    analysis = ReferenceUnderstandingResult(
        summary=str(analysis_raw.get("summary", "") if isinstance(analysis_raw, dict) else ""),
        dominant_themes=[
            str(item).strip()
            for item in analysis_raw.get("dominant_themes", [])
            if str(item).strip()
        ]
        if isinstance(analysis_raw, dict) and isinstance(analysis_raw.get("dominant_themes"), list)
        else [],
        domain_terms=[
            str(item).strip()
            for item in analysis_raw.get("domain_terms", [])
            if str(item).strip()
        ]
        if isinstance(analysis_raw, dict) and isinstance(analysis_raw.get("domain_terms"), list)
        else [],
        document_patterns=[
            str(item).strip()
            for item in analysis_raw.get("document_patterns", [])
            if str(item).strip()
        ]
        if isinstance(analysis_raw, dict) and isinstance(analysis_raw.get("document_patterns"), list)
        else [],
        signal_scores={
            str(key): float(value)
            for key, value in analysis_raw.get("signal_scores", {}).items()
        }
        if isinstance(analysis_raw, dict) and isinstance(analysis_raw.get("signal_scores"), dict)
        else {},
        low_signal=bool(analysis_raw.get("low_signal", False))
        if isinstance(analysis_raw, dict)
        else False,
    )
    return NormalizedDocument(
        document_summary=str(payload.get("document_summary", "") if isinstance(payload, dict) else ""),
        sections=sections,
        analysis=analysis,
    )


class DocxReferenceExtractor:
    def extract(self, payload: bytes) -> NormalizedDocument:
        document = DocxDocument(BytesIO(payload))
        groups: list[tuple[str, list[str], dict[str, object]]] = []
        current_title = "Dokument"
        current_lines: list[str] = []
        ordinal = 1

        def flush(locator: dict[str, object]) -> None:
            nonlocal ordinal, current_lines, current_title
            text = "\n".join(line for line in current_lines if line.strip()).strip()
            if not text:
                return
            groups.append((current_title, [text], locator))
            ordinal += 1

        para_index = 0
        for paragraph in document.paragraphs:
            para_index += 1
            text = (paragraph.text or "").strip()
            if not text:
                continue
            style_name = ((paragraph.style.name if paragraph.style is not None else "") or "").lower()
            if style_name.startswith("heading") and current_lines:
                flush({"paragraph": para_index - 1})
                current_lines = []
            if style_name.startswith("heading"):
                current_title = text
                continue
            current_lines.append(text)

        if current_lines:
            flush({"paragraph": para_index})

        table_index = 0
        for table in document.tables:
            table_index += 1
            rows: list[str] = []
            for row in table.rows:
                row_values = [cell.text.strip() for cell in row.cells if (cell.text or "").strip()]
                if row_values:
                    rows.append(" | ".join(row_values))
            if rows:
                groups.append((f"Tabelle {table_index}", rows, {"table": table_index}))

        return _build_document_from_groups(groups, default_title="DOCX Referenz")


class PptxReferenceExtractor:
    def extract(self, payload: bytes) -> NormalizedDocument:
        presentation = Presentation(BytesIO(payload))
        groups: list[tuple[str, list[str], dict[str, object]]] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            lines: list[str] = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text = (shape.text or "").strip()
                    if text:
                        lines.extend(line.strip() for line in text.splitlines() if line.strip())
                    continue
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        row_values = [cell.text.strip() for cell in row.cells if (cell.text or "").strip()]
                        if row_values:
                            lines.append(" | ".join(row_values))
            if not lines:
                continue
            title = lines[0][:120] if lines else f"Folie {slide_index}"
            groups.append((title, lines, {"slide": slide_index}))
        return _build_document_from_groups(groups, default_title="PPTX Referenz")


class XlsxReferenceExtractor:
    def extract(self, payload: bytes) -> NormalizedDocument:
        workbook = load_workbook(BytesIO(payload), data_only=True, read_only=True)
        groups: list[tuple[str, list[str], dict[str, object]]] = []
        for sheet in workbook.worksheets:
            row_lines: list[tuple[int, str]] = []
            for row_idx, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                values = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
                if values:
                    row_lines.append((row_idx, " | ".join(values)))
            if not row_lines:
                continue
            chunk: list[str] = []
            chunk_start = row_lines[0][0]
            current_len = 0
            for row_idx, line in row_lines:
                extra_len = len(line) + 1
                if chunk and current_len + extra_len > MAX_SECTION_CHARS:
                    groups.append(
                        (
                            f"{sheet.title} ({chunk_start}-{row_idx - 1})",
                            chunk,
                            {"sheet": sheet.title, "row_start": chunk_start, "row_end": row_idx - 1},
                        )
                    )
                    chunk = []
                    chunk_start = row_idx
                    current_len = 0
                chunk.append(f"Zeile {row_idx}: {line}")
                current_len += extra_len
            if chunk:
                groups.append(
                    (
                        f"{sheet.title} ({chunk_start}-{row_lines[-1][0]})",
                        chunk,
                        {"sheet": sheet.title, "row_start": chunk_start, "row_end": row_lines[-1][0]},
                    )
                )
        return _build_document_from_groups(groups, default_title="XLSX Referenz")


class TextReferenceExtractor:
    def __init__(self, *, markdown: bool = False) -> None:
        self.markdown = markdown

    def extract(self, payload: bytes) -> NormalizedDocument:
        text = payload.decode("utf-8", errors="ignore").replace("\r\n", "\n").strip()
        if not text:
            raise ReferenceExtractionError("Text content is empty")
        groups = _extract_markdown_groups(text) if self.markdown else _extract_plaintext_groups(text)
        return _build_document_from_groups(groups, default_title="Text Referenz")


class PdfReferenceExtractor:
    def extract(self, payload: bytes) -> NormalizedDocument:
        try:
            reader = PdfReader(BytesIO(payload))
        except Exception as exc:  # noqa: BLE001
            raise ReferenceExtractionError(f"PDF extraction failed: {exc}") from exc

        groups: list[tuple[str, list[str], dict[str, object]]] = []
        for page_index, page in enumerate(reader.pages, start=1):
            page_text = (page.extract_text() or "").strip()
            if not page_text:
                continue
            paragraphs = [item.strip() for item in re.split(r"\n\s*\n+", page_text) if item.strip()]
            groups.append((f"Seite {page_index}", paragraphs or [page_text], {"page": page_index}))
        if not groups:
            raise ReferenceExtractionError("PDF contains no extractable text")
        return _build_document_from_groups(groups, default_title="PDF Referenz")


def get_reference_extractor(file_type: str) -> ReferenceExtractor:
    if file_type == "DOCX":
        return DocxReferenceExtractor()
    if file_type == "PPTX":
        return PptxReferenceExtractor()
    if file_type == "XLSX":
        return XlsxReferenceExtractor()
    if file_type == "TXT":
        return TextReferenceExtractor(markdown=False)
    if file_type == "MD":
        return TextReferenceExtractor(markdown=True)
    if file_type == "PDF":
        return PdfReferenceExtractor()
    raise ReferenceExtractionError("Unsupported reference file type")


def extract_reference_document(*, payload: bytes, file_type: str) -> NormalizedDocument:
    extractor = get_reference_extractor(file_type)
    return extractor.extract(payload)


def _extract_markdown_groups(text: str) -> list[tuple[str, list[str], dict[str, object]]]:
    groups: list[tuple[str, list[str], dict[str, object]]] = []
    current_title = "Dokument"
    current_lines: list[str] = []
    heading_re = re.compile(r"^(#{1,6})\s+(.*)$")

    def flush(line_no: int) -> None:
        nonlocal current_lines
        content = "\n".join(current_lines).strip()
        if content:
            groups.append((current_title, [content], {"line": line_no}))
        current_lines = []

    for line_no, raw_line in enumerate(text.splitlines(), start=1):
        line = raw_line.rstrip()
        match = heading_re.match(line)
        if match:
            flush(line_no - 1)
            current_title = match.group(2).strip() or current_title
            continue
        if line.strip():
            current_lines.append(line.strip())
        elif current_lines:
            current_lines.append("")
    flush(len(text.splitlines()))
    return groups


def _extract_plaintext_groups(text: str) -> list[tuple[str, list[str], dict[str, object]]]:
    groups: list[tuple[str, list[str], dict[str, object]]] = []
    for index, chunk in enumerate(re.split(r"\n\s*\n+", text), start=1):
        cleaned = chunk.strip()
        if not cleaned:
            continue
        title = cleaned.splitlines()[0][:120]
        groups.append((title or f"Abschnitt {index}", [cleaned], {"group": index}))
    return groups


def _build_document_from_groups(
    groups: list[tuple[str, list[str], dict[str, object]]],
    *,
    default_title: str,
) -> NormalizedDocument:
    sections: list[NormalizedSection] = []
    ordinal = 1
    for group_title, lines, locator in groups:
        buffer: list[str] = []
        current_len = 0
        part = 1
        for raw_line in lines:
            line = raw_line.strip()
            if not line:
                if buffer and buffer[-1] != "":
                    buffer.append("")
                continue
            if buffer and current_len + len(line) + 1 > MAX_SECTION_CHARS:
                content = "\n".join(buffer).strip()
                if content:
                    title = group_title if part == 1 else f"{group_title} ({part})"
                    section_kind, keywords, themes, signals = _analyze_section(
                        title=title,
                        content=content,
                        locator={**locator, "part": part},
                    )
                    sections.append(
                        NormalizedSection(
                            id=f"section-{ordinal}",
                            type="section",
                            title=title,
                            locator={**locator, "part": part},
                            content=content,
                            estimated_tokens=estimate_token_count(content),
                            section_kind=section_kind,
                            keywords=keywords,
                            themes=themes,
                            signals=signals,
                        )
                    )
                    ordinal += 1
                    part += 1
                buffer = []
                current_len = 0
            buffer.append(line)
            current_len += len(line) + 1
        content = "\n".join(buffer).strip()
        if content:
            title = group_title or default_title
            if part > 1:
                title = f"{title} ({part})"
            section_kind, keywords, themes, signals = _analyze_section(
                title=title,
                content=content,
                locator={**locator, "part": part},
            )
            sections.append(
                NormalizedSection(
                    id=f"section-{ordinal}",
                    type="section",
                    title=title,
                    locator={**locator, "part": part},
                    content=content,
                    estimated_tokens=estimate_token_count(content),
                    section_kind=section_kind,
                    keywords=keywords,
                    themes=themes,
                    signals=signals,
                )
            )
            ordinal += 1

    if not sections:
        raise ReferenceExtractionError("No extractable sections found")

    analysis = _analyze_document(sections=sections, default_title=default_title)
    summary_parts: list[str] = []
    if analysis.dominant_themes:
        summary_parts.append(
            f"Themen: {', '.join(analysis.dominant_themes[:4])}"
        )
    if analysis.domain_terms:
        summary_parts.append(
            f"Begriffe: {', '.join(analysis.domain_terms[:6])}"
        )
    for section in sections[:2]:
        header = section.title or default_title
        excerpt = section.content[:140].strip()
        if excerpt:
            summary_parts.append(f"{header}: {excerpt}")
    summary = "\n".join(summary_parts)[:SUMMARY_MAX_CHARS].strip() or analysis.summary
    return NormalizedDocument(document_summary=summary, sections=sections, analysis=analysis)


def _tokenize_terms(text: str) -> list[str]:
    return [
        term
        for term in re.findall(r"[A-Za-zÄÖÜäöüß][A-Za-zÄÖÜäöüß0-9_-]{2,}", text.lower())
        if term not in STOPWORDS
    ]


def _extract_keywords(text: str, *, limit: int = 10) -> list[str]:
    counter = Counter(_tokenize_terms(text))
    return [term for term, _count in counter.most_common(limit)]


def _score_theme_matches(text: str) -> dict[str, float]:
    haystack = text.lower()
    scores: dict[str, float] = {}
    for theme, patterns in THEME_PATTERNS.items():
        score = 0.0
        for pattern in patterns:
            occurrences = haystack.count(pattern)
            if occurrences:
                score += 1.0 + math.log1p(occurrences)
        if score > 0:
            scores[theme] = round(score, 3)
    return scores


def _analyze_section(
    *,
    title: str,
    content: str,
    locator: dict[str, object],
) -> tuple[str, list[str], list[str], dict[str, float]]:
    text = f"{title}\n{content}".strip()
    keywords = _extract_keywords(text)
    signals = _score_theme_matches(text)
    themes = [
        theme
        for theme, score in sorted(signals.items(), key=lambda item: item[1], reverse=True)
        if score > 0
    ]

    if "table" in locator:
        section_kind = "table"
    elif "sheet" in locator:
        section_kind = "table"
    elif re.search(r"(^|\n)\s*(?:[-*•]|\d+[.)])\s+", content):
        section_kind = "bullet_list"
    elif themes and themes[0] in {"workflow", "instruction"}:
        section_kind = "procedure"
    elif title and len(title.split()) <= 12:
        section_kind = "heading_block"
    else:
        section_kind = "section"

    return section_kind, keywords[:8], themes[:4], signals


def _analyze_document(
    *,
    sections: list[NormalizedSection],
    default_title: str,
) -> ReferenceUnderstandingResult:
    aggregate_scores: Counter[str] = Counter()
    keyword_counter: Counter[str] = Counter()
    for section in sections:
        aggregate_scores.update(section.signals)
        keyword_counter.update(section.keywords)

    dominant_themes = [
        theme
        for theme, score in aggregate_scores.most_common()
        if float(score) > 0
    ][:4]
    domain_terms = [term for term, _count in keyword_counter.most_common(12)]
    summary = (
        f"{default_title}: "
        f"{', '.join(dominant_themes[:3]) or 'allgemeine Referenz'}"
    )
    low_signal = sum(section.estimated_tokens for section in sections) < 80 or not dominant_themes
    return ReferenceUnderstandingResult(
        summary=summary[:SUMMARY_MAX_CHARS],
        dominant_themes=dominant_themes,
        domain_terms=domain_terms[:10],
        document_patterns=dominant_themes[:4],
        signal_scores={theme: round(float(score), 3) for theme, score in aggregate_scores.items()},
        low_signal=low_signal,
    )
