from __future__ import annotations

from io import BytesIO

from .asset_placeholders import collect_asset_occurrences
from .asset_resolver import asset_missing_error
from .image_utils import resolve_display_size
from .office_asset_types import ResolvedOfficeAsset


def inject_pptx_assets(
    *,
    payload: bytes,
    aliases_by_key: dict[str, list[str]],
    assets_by_key: dict[str, ResolvedOfficeAsset | None],
    fail_on_missing_asset: bool,
    error_by_key: dict[str, dict[str, object]] | None = None,
) -> tuple[bytes, list[dict[str, object]], dict[str, int]]:
    from pptx import Presentation  # type: ignore

    presentation = Presentation(BytesIO(payload))
    errors: list[dict[str, object]] = []
    occurrences: dict[str, int] = {}
    available_keys = set(aliases_by_key.keys())
    per_key_errors = error_by_key or {}

    for slide in presentation.slides:
        for shape in list(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                continue

            text = shape.text or ""
            if not text:
                continue

            matches = collect_asset_occurrences(text, available_keys)
            if not matches:
                continue

            left = int(shape.left)
            top = int(shape.top)
            width = int(shape.width)
            height = int(shape.height)
            updated = text

            for match in matches:
                key = match.key
                occurrences[key] = occurrences.get(key, 0) + 1
                resolved = assets_by_key.get(key)
                if resolved is None:
                    if fail_on_missing_asset:
                        errors.append(per_key_errors.get(key) or asset_missing_error(key))
                    continue

                img_w, img_h = resolve_display_size(
                    source_width=resolved.width,
                    source_height=resolved.height,
                    request_width=match.width_px,
                    request_height=match.height_px,
                    max_width=width,
                    max_height=height,
                )
                x = left + max(0, (width - img_w) // 2)
                y = top + max(0, (height - img_h) // 2)
                slide.shapes.add_picture(BytesIO(resolved.payload), x, y, width=img_w, height=img_h)
                updated = updated.replace(match.raw, "")

            if fail_on_missing_asset and any(assets_by_key.get(match.key) is None for match in matches):
                continue
            shape.text = updated.strip()

    out = BytesIO()
    presentation.save(out)
    return out.getvalue(), errors, occurrences
