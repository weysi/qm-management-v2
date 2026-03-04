from __future__ import annotations

from io import BytesIO
from pathlib import Path
import tempfile
from unittest.mock import patch
from zipfile import ZIP_DEFLATED, ZipFile
import base64
import hashlib

from django.core.files.uploadedfile import SimpleUploadedFile
from django.test import TestCase, override_settings
from rest_framework.test import APIClient
from docx import Document as WordDocument

from documents.models import Document, DocumentVersion, WorkspaceAsset


class DocumentApiTests(TestCase):
    VALID_PNG = base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO0N8v8AAAAASUVORK5CYII="
    )

    def setUp(self):
        self.client = APIClient()
        self.temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(self.temp_dir.cleanup)

    @override_settings(DOCUMENTS_DATA_ROOT=Path("/tmp"))
    def test_upload_and_tree_and_soft_delete(self):
        payload = b"Hello {{user.name}} and {{assets.logo}}"
        uploaded = SimpleUploadedFile("sample.txt", payload, content_type="text/plain")
        res = self.client.post(
            "/api/v1/documents/upload",
            {"handbook_id": "hb-1", "file": uploaded, "path": "folder/sample.txt"},
            format="multipart",
        )
        self.assertEqual(res.status_code, 201)
        body = res.json()
        self.assertEqual(body["kind"], "file")
        doc_id = body["document"]["id"]

        tree = self.client.get("/api/v1/files/tree", {"handbook_id": "hb-1"})
        self.assertEqual(tree.status_code, 200)
        self.assertTrue(tree.json()["tree"])

        delete = self.client.delete(
            "/api/v1/files",
            {"handbook_id": "hb-1", "path": "folder/sample.txt", "recursive": False},
            format="json",
        )
        self.assertEqual(delete.status_code, 200)
        self.assertEqual(Document.objects.filter(id=doc_id, deleted_at__isnull=False).count(), 1)

    @override_settings(DOCUMENTS_DATA_ROOT=Path("/tmp"))
    def test_upload_links_to_handbook_and_list_returns_document(self):
        handbook_id = "hb-link"
        uploaded = SimpleUploadedFile("linked.txt", b"Hello {{company.name}}", content_type="text/plain")
        upload_res = self.client.post(
            "/api/v1/documents/upload",
            {"handbook_id": handbook_id, "file": uploaded, "path": "linked.txt"},
            format="multipart",
        )
        self.assertEqual(upload_res.status_code, 201)
        payload = upload_res.json()
        document_id = payload["document"]["id"]
        self.assertEqual(payload["document"]["handbook_id"], handbook_id)

        doc = Document.objects.get(id=document_id)
        self.assertEqual(doc.handbook_id, handbook_id)

        list_res = self.client.get("/api/v1/documents", {"handbook_id": handbook_id})
        self.assertEqual(list_res.status_code, 200)
        listed_ids = [item["id"] for item in list_res.json()["documents"]]
        self.assertIn(document_id, listed_ids)

    @override_settings(DOCUMENTS_DATA_ROOT=Path("/tmp"))
    def test_files_tree_contract_for_file_nodes(self):
        handbook_id = "hb-tree"

        upload_one = SimpleUploadedFile("a.txt", b"Hello", content_type="text/plain")
        upload_two = SimpleUploadedFile("b.txt", b"World", content_type="text/plain")
        self.client.post(
            "/api/v1/documents/upload",
            {"handbook_id": handbook_id, "file": upload_one, "path": "a.txt"},
            format="multipart",
        )
        self.client.post(
            "/api/v1/documents/upload",
            {"handbook_id": handbook_id, "file": upload_two, "path": "folder/b.txt"},
            format="multipart",
        )

        tree_res = self.client.get("/api/v1/files/tree", {"handbook_id": handbook_id})
        self.assertEqual(tree_res.status_code, 200)
        tree = tree_res.json()["tree"]
        self.assertTrue(tree)

        def walk(nodes):
            for node in nodes:
                yield node
                children = node.get("children")
                if isinstance(children, list):
                    yield from walk(children)

        flattened = list(walk(tree))
        folder_nodes = [node for node in flattened if node["kind"] == "folder"]
        file_nodes = [node for node in flattened if node["kind"] == "file"]

        self.assertTrue(folder_nodes)
        self.assertTrue(file_nodes)
        self.assertTrue(all(isinstance(node.get("children"), list) for node in folder_nodes))
        self.assertTrue(all("children" not in node for node in file_nodes))

    @override_settings(DOCUMENTS_DATA_ROOT=Path("/tmp"))
    def test_logo_upload_returns_asset_metadata_and_download(self):
        handbook_id = "hb-assets"
        logo_bytes = b"\x89PNG\r\n\x1a\n\x00\x00\x00\x00"
        upload = SimpleUploadedFile("logo.png", logo_bytes, content_type="image/png")
        res = self.client.post(
            f"/api/v1/handbooks/{handbook_id}/assets",
            {"file": upload, "asset_type": "logo"},
            format="multipart",
        )
        self.assertEqual(res.status_code, 201)
        asset = res.json()["asset"]
        self.assertEqual(asset["asset_type"], "logo")
        self.assertEqual(asset["filename"], "logo.png")
        self.assertEqual(asset["status"], "READY")
        self.assertTrue(asset["download_url"].endswith(f"/api/v1/handbooks/{handbook_id}/assets/logo/download"))
        self.assertEqual(asset["preview_url"], asset["download_url"])

        listed = self.client.get(f"/api/v1/handbooks/{handbook_id}/assets")
        self.assertEqual(listed.status_code, 200)
        payload = listed.json()["assets"]
        self.assertEqual(len(payload), 1)
        self.assertEqual(payload[0]["asset_type"], "logo")

        download = self.client.get(f"/api/v1/handbooks/{handbook_id}/assets/logo/download")
        self.assertEqual(download.status_code, 200)
        self.assertEqual(b"".join(download.streaming_content), logo_bytes)

        delete = self.client.delete(f"/api/v1/handbooks/{handbook_id}/assets/logo")
        self.assertEqual(delete.status_code, 200)

    @override_settings(DOCUMENTS_DATA_ROOT=Path("/tmp"))
    def test_zip_upload_auto_binds_logo_and_signature(self):
        handbook_id = "hb-zip"
        archive = BytesIO()
        with ZipFile(archive, "w", compression=ZIP_DEFLATED) as bundle:
            bundle.writestr("templates/policy.txt", "Policy {{company.name}}")
            bundle.writestr("logo.png", b"\x89PNG\r\n")
            bundle.writestr("signature.png", b"\x89PNG\r\nSIG")
            bundle.writestr("notes.tmp", "skip me")

        upload = SimpleUploadedFile("bundle.zip", archive.getvalue(), content_type="application/zip")
        res = self.client.post(
            "/api/v1/documents/upload",
            {"handbook_id": handbook_id, "file": upload},
            format="multipart",
        )
        self.assertEqual(res.status_code, 201)
        body = res.json()
        self.assertEqual(body["kind"], "zip")
        self.assertEqual(body["summary"]["documents_created"], 1)
        self.assertEqual(body["summary"]["assets_bound"], 2)
        self.assertGreaterEqual(body["summary"]["warnings"], 1)

        listed_docs = self.client.get("/api/v1/documents", {"handbook_id": handbook_id})
        self.assertEqual(listed_docs.status_code, 200)
        self.assertEqual(len(listed_docs.json()["documents"]), 1)

        listed_assets = self.client.get(f"/api/v1/handbooks/{handbook_id}/assets")
        self.assertEqual(listed_assets.status_code, 200)
        asset_types = sorted(item["asset_type"] for item in listed_assets.json()["assets"])
        self.assertEqual(asset_types, ["logo", "signature"])

    @override_settings(DOCUMENTS_DATA_ROOT=Path("/tmp"))
    def test_render_resolves_assets_for_mustache_and_legacy_alias(self):
        handbook_id = "hb-render-assets"
        logo_upload = SimpleUploadedFile("logo.png", b"\x89PNG\r\n", content_type="image/png")
        logo_res = self.client.post(
            f"/api/v1/handbooks/{handbook_id}/assets",
            {"file": logo_upload, "asset_type": "logo"},
            format="multipart",
        )
        self.assertEqual(logo_res.status_code, 201)

        template = SimpleUploadedFile(
            "template.txt",
            b"mustache={{assets.logo}}\nlegacy=[LOGO]\n",
            content_type="text/plain",
        )
        upload_res = self.client.post(
            "/api/v1/documents/upload",
            {"handbook_id": handbook_id, "file": template},
            format="multipart",
        )
        self.assertEqual(upload_res.status_code, 201)
        document_id = upload_res.json()["document"]["id"]

        render_res = self.client.post(
            f"/api/v1/documents/{document_id}/render",
            {"variables": {}},
            format="json",
        )
        self.assertEqual(render_res.status_code, 200)
        version_id = render_res.json()["version"]["id"]
        version = DocumentVersion.objects.get(id=version_id)
        output = Path(version.file_path).read_text(encoding="utf-8")

        expected = f"/api/v1/handbooks/{handbook_id}/assets/logo/download"
        self.assertIn(expected, output)
        self.assertNotIn("[LOGO]", output)

    @override_settings(DOCUMENTS_DATA_ROOT=Path("/tmp"))
    def test_docx_render_embeds_logo_image_instead_of_url_text(self):
        handbook_id = "hb-docx-logo"
        logo_upload = SimpleUploadedFile("logo.png", self.VALID_PNG, content_type="image/png")
        logo_res = self.client.post(
            f"/api/v1/handbooks/{handbook_id}/assets",
            {"file": logo_upload, "asset_type": "logo"},
            format="multipart",
        )
        self.assertEqual(logo_res.status_code, 201)

        doc = WordDocument()
        doc.add_paragraph("Header")
        doc.add_paragraph("{{assets.logo}}")
        doc_stream = BytesIO()
        doc.save(doc_stream)
        upload = SimpleUploadedFile(
            "template.docx",
            doc_stream.getvalue(),
            content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )
        upload_res = self.client.post(
            "/api/v1/documents/upload",
            {"handbook_id": handbook_id, "file": upload},
            format="multipart",
        )
        self.assertEqual(upload_res.status_code, 201)
        document_id = upload_res.json()["document"]["id"]

        render_res = self.client.post(
            f"/api/v1/documents/{document_id}/render",
            {"variables": {}},
            format="json",
        )
        self.assertEqual(render_res.status_code, 200)
        version_id = render_res.json()["version"]["id"]
        version = DocumentVersion.objects.get(id=version_id)

        with ZipFile(Path(version.file_path), "r") as archive:
            media_entries = [name for name in archive.namelist() if name.startswith("word/media/")]
            self.assertTrue(media_entries)
            media_bytes = archive.read(media_entries[0])
            self.assertEqual(
                hashlib.sha256(media_bytes).hexdigest(),
                hashlib.sha256(self.VALID_PNG).hexdigest(),
            )
            xml = archive.read("word/document.xml").decode("utf-8", errors="ignore")
            self.assertNotIn("/api/v1/handbooks/", xml)
            self.assertNotIn("__ASSET_LOGO__", xml)

    @override_settings(DOCUMENTS_DATA_ROOT=Path("/tmp"))
    def test_docx_missing_asset_fails_by_default_policy(self):
        handbook_id = "hb-docx-fail"
        doc = WordDocument()
        doc.add_paragraph("__ASSET_LOGO__")
        stream = BytesIO()
        doc.save(stream)
        upload = SimpleUploadedFile(
            "missing-logo.docx",
            stream.getvalue(),
            content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )
        upload_res = self.client.post(
            "/api/v1/documents/upload",
            {"handbook_id": handbook_id, "file": upload},
            format="multipart",
        )
        self.assertEqual(upload_res.status_code, 201)
        document_id = upload_res.json()["document"]["id"]

        render_res = self.client.post(
            f"/api/v1/documents/{document_id}/render",
            {"variables": {}},
            format="json",
        )
        self.assertEqual(render_res.status_code, 400)
        errors = render_res.json().get("errors", [])
        self.assertTrue(any(err.get("error_code") == "MISSING_REQUIRED_ASSET" for err in errors))

    @override_settings(DOCUMENTS_DATA_ROOT=Path("/tmp"))
    def test_docx_keep_placeholder_policy(self):
        handbook_id = "hb-docx-keep"
        doc = WordDocument()
        doc.add_paragraph("__ASSET_LOGO__")
        stream = BytesIO()
        doc.save(stream)
        upload = SimpleUploadedFile(
            "keep-logo.docx",
            stream.getvalue(),
            content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )
        upload_res = self.client.post(
            "/api/v1/documents/upload",
            {"handbook_id": handbook_id, "file": upload},
            format="multipart",
        )
        self.assertEqual(upload_res.status_code, 201)
        document_id = upload_res.json()["document"]["id"]

        render_res = self.client.post(
            f"/api/v1/documents/{document_id}/render",
            {"variables": {}, "generation_policy": {"on_missing_asset": "KEEP_PLACEHOLDER"}},
            format="json",
        )
        self.assertEqual(render_res.status_code, 200)
        version_id = render_res.json()["version"]["id"]
        version = DocumentVersion.objects.get(id=version_id)
        with ZipFile(Path(version.file_path), "r") as archive:
            xml = archive.read("word/document.xml").decode("utf-8", errors="ignore")
            self.assertIn("assets.logo", xml)

    @override_settings(DOCUMENTS_DATA_ROOT=Path("/tmp"))
    def test_pptx_render_replaces_logo_placeholder_with_picture(self):
        from pptx import Presentation  # type: ignore
        from pptx.util import Inches  # type: ignore

        handbook_id = "hb-pptx-logo"
        logo_upload = SimpleUploadedFile("logo.png", self.VALID_PNG, content_type="image/png")
        self.client.post(
            f"/api/v1/handbooks/{handbook_id}/assets",
            {"file": logo_upload, "asset_type": "logo"},
            format="multipart",
        )

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        box.text_frame.text = "[LOGO]"
        stream = BytesIO()
        presentation.save(stream)

        upload = SimpleUploadedFile(
            "template.pptx",
            stream.getvalue(),
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        upload_res = self.client.post(
            "/api/v1/documents/upload",
            {"handbook_id": handbook_id, "file": upload},
            format="multipart",
        )
        self.assertEqual(upload_res.status_code, 201)
        document_id = upload_res.json()["document"]["id"]

        render_res = self.client.post(
            f"/api/v1/documents/{document_id}/render",
            {"variables": {}},
            format="json",
        )
        self.assertEqual(render_res.status_code, 200)
        version = DocumentVersion.objects.get(id=render_res.json()["version"]["id"])

        with ZipFile(Path(version.file_path), "r") as archive:
            media_entries = [name for name in archive.namelist() if name.startswith("ppt/media/")]
            self.assertTrue(media_entries)
            media_bytes = archive.read(media_entries[0])
            self.assertEqual(
                hashlib.sha256(media_bytes).hexdigest(),
                hashlib.sha256(self.VALID_PNG).hexdigest(),
            )
            slide_xml = archive.read("ppt/slides/slide1.xml").decode("utf-8", errors="ignore")
            self.assertNotIn("LOGO", slide_xml)
            self.assertNotIn("assets.logo", slide_xml)
            self.assertNotIn("/api/v1/handbooks/", slide_xml)

    @override_settings(DOCUMENTS_DATA_ROOT=Path("/tmp"))
    def test_xlsx_render_replaces_logo_placeholder_with_anchored_image(self):
        from openpyxl import Workbook, load_workbook  # type: ignore

        handbook_id = "hb-xlsx-logo"
        logo_upload = SimpleUploadedFile("logo.png", self.VALID_PNG, content_type="image/png")
        self.client.post(
            f"/api/v1/handbooks/{handbook_id}/assets",
            {"file": logo_upload, "asset_type": "logo"},
            format="multipart",
        )

        workbook = Workbook()
        ws = workbook.active
        ws["A1"] = "__ASSET_LOGO__"
        stream = BytesIO()
        workbook.save(stream)

        upload = SimpleUploadedFile(
            "template.xlsx",
            stream.getvalue(),
            content_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )
        upload_res = self.client.post(
            "/api/v1/documents/upload",
            {"handbook_id": handbook_id, "file": upload},
            format="multipart",
        )
        self.assertEqual(upload_res.status_code, 201)
        document_id = upload_res.json()["document"]["id"]

        render_res = self.client.post(
            f"/api/v1/documents/{document_id}/render",
            {"variables": {}},
            format="json",
        )
        self.assertEqual(render_res.status_code, 200)
        version = DocumentVersion.objects.get(id=render_res.json()["version"]["id"])
        wb = load_workbook(version.file_path)
        self.assertIn(len(wb.active._images), {1, 2})
        self.assertIn(wb.active["A1"].value, {None, ""})
        with ZipFile(Path(version.file_path), "r") as archive:
            media_entries = [name for name in archive.namelist() if name.startswith("xl/media/")]
            self.assertTrue(media_entries)
            media_bytes = archive.read(media_entries[0])
            self.assertEqual(
                hashlib.sha256(media_bytes).hexdigest(),
                hashlib.sha256(self.VALID_PNG).hexdigest(),
            )

    @override_settings(DOCUMENTS_DATA_ROOT=Path("/tmp"))
    def test_render_fails_when_asset_bytes_do_not_match_stored_hash(self):
        handbook_id = "hb-integrity"
        logo_upload = SimpleUploadedFile("logo.png", self.VALID_PNG, content_type="image/png")
        upload_res = self.client.post(
            f"/api/v1/handbooks/{handbook_id}/assets",
            {"file": logo_upload, "asset_type": "logo"},
            format="multipart",
        )
        self.assertEqual(upload_res.status_code, 201)

        asset = WorkspaceAsset.objects.filter(
            handbook_id=handbook_id,
            asset_type="logo",
            deleted_at__isnull=True,
        ).first()
        self.assertIsNotNone(asset)
        Path(str(asset.file_path)).write_bytes(b"tampered-content")

        doc = WordDocument()
        doc.add_paragraph("{{assets.logo}}")
        stream = BytesIO()
        doc.save(stream)
        template_upload = SimpleUploadedFile(
            "integrity.docx",
            stream.getvalue(),
            content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )
        create_res = self.client.post(
            "/api/v1/documents/upload",
            {"handbook_id": handbook_id, "file": template_upload},
            format="multipart",
        )
        self.assertEqual(create_res.status_code, 201)
        document_id = create_res.json()["document"]["id"]

        render_res = self.client.post(
            f"/api/v1/documents/{document_id}/render",
            {"variables": {}},
            format="json",
        )
        self.assertEqual(render_res.status_code, 400)
        errors = render_res.json().get("errors", [])
        self.assertTrue(any(err.get("error_code") == "ASSET_INTEGRITY_ERROR" for err in errors))

    @override_settings(DOCUMENTS_DATA_ROOT=Path("/tmp"))
    def test_docx_render_preserves_non_placeholder_run_styling(self):
        handbook_id = "hb-docx-run-style"
        logo_upload = SimpleUploadedFile("logo.png", self.VALID_PNG, content_type="image/png")
        self.client.post(
            f"/api/v1/handbooks/{handbook_id}/assets",
            {"file": logo_upload, "asset_type": "logo"},
            format="multipart",
        )

        doc = WordDocument()
        paragraph = doc.add_paragraph()
        before = paragraph.add_run("Before ")
        before.bold = True
        paragraph.add_run("{{assets.logo}}")
        after = paragraph.add_run(" After")
        after.italic = True
        stream = BytesIO()
        doc.save(stream)

        upload = SimpleUploadedFile(
            "style.docx",
            stream.getvalue(),
            content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )
        upload_res = self.client.post(
            "/api/v1/documents/upload",
            {"handbook_id": handbook_id, "file": upload},
            format="multipart",
        )
        self.assertEqual(upload_res.status_code, 201)

        document_id = upload_res.json()["document"]["id"]
        render_res = self.client.post(
            f"/api/v1/documents/{document_id}/render",
            {"variables": {}},
            format="json",
        )
        self.assertEqual(render_res.status_code, 200)
        version = DocumentVersion.objects.get(id=render_res.json()["version"]["id"])
        rendered = WordDocument(version.file_path)
        rendered_paragraph = rendered.paragraphs[0]
        self.assertIn("Before", rendered_paragraph.text)
        self.assertIn("After", rendered_paragraph.text)
        self.assertNotIn("assets.logo", rendered_paragraph.text)
        self.assertTrue(any((run.text or "").startswith("Before") and run.bold for run in rendered_paragraph.runs))
        self.assertTrue(any((run.text or "").endswith("After") and run.italic for run in rendered_paragraph.runs))

    @override_settings(DOCUMENTS_DATA_ROOT=Path("/tmp"))
    @patch("documents.services.variable_fill_service.AiClient.generate_variable_value")
    def test_ai_fill_variable_success(self, mock_generate):
        from documents.services.ai_client import VariableValueResponse

        mock_generate.return_value = VariableValueResponse(
            value="ISO 9001 konformer Beispieltext",
            model="gpt-4o-mini",
            usage={"prompt_tokens": 10, "completion_tokens": 8, "total_tokens": 18},
        )

        res = self.client.post(
            "/api/v1/handbooks/hb-ai/variables/ai-fill",
            {
                "variable_name": "company.scope",
                "current_value": "",
                "instruction": "Kurz und präzise formulieren.",
                "client_context": {"company_name": "Beispiel GmbH"},
                "language": "de-DE",
                "constraints": {"max_length": 120, "required": True},
            },
            format="json",
        )
        self.assertEqual(res.status_code, 200)
        body = res.json()
        self.assertEqual(body["value"], "ISO 9001 konformer Beispieltext")
        self.assertEqual(body["model"], "gpt-4o-mini")
        self.assertEqual(body["usage"]["total_tokens"], 18)

    @override_settings(DOCUMENTS_DATA_ROOT=Path("/tmp"))
    def test_ai_fill_variable_requires_instruction(self):
        res = self.client.post(
            "/api/v1/handbooks/hb-ai/variables/ai-fill",
            {
                "variable_name": "company.scope",
                "instruction": "",
                "client_context": {},
                "language": "de-DE",
                "constraints": {"max_length": 120, "required": True},
            },
            format="json",
        )
        self.assertEqual(res.status_code, 400)
        self.assertEqual(res.json().get("error_code"), "INSTRUCTION_REQUIRED")

    @override_settings(DOCUMENTS_DATA_ROOT=Path("/tmp"))
    @patch("documents.services.variable_fill_service.AiClient.generate_variable_value")
    def test_ai_fill_variable_max_length_hard_fail(self, mock_generate):
        from documents.services.ai_client import VariableValueResponse

        mock_generate.return_value = VariableValueResponse(
            value="x" * 25,
            model="gpt-4o-mini",
            usage={"prompt_tokens": 1, "completion_tokens": 1, "total_tokens": 2},
        )

        res = self.client.post(
            "/api/v1/handbooks/hb-ai/variables/ai-fill",
            {
                "variable_name": "company.scope",
                "instruction": "Kurz",
                "client_context": {},
                "language": "de-DE",
                "constraints": {"max_length": 10, "required": True},
            },
            format="json",
        )
        self.assertEqual(res.status_code, 400)
        self.assertEqual(res.json().get("error_code"), "MAX_LENGTH_EXCEEDED")

    @override_settings(DOCUMENTS_DATA_ROOT=Path("/tmp"))
    @patch("documents.services.variable_fill_service.AiClient.generate_variable_value")
    def test_ai_fill_variable_maps_provider_error_to_502(self, mock_generate):
        from documents.services.ai_client import AiClientError

        mock_generate.side_effect = AiClientError("upstream down")
        res = self.client.post(
            "/api/v1/handbooks/hb-ai/variables/ai-fill",
            {
                "variable_name": "company.scope",
                "instruction": "Kurz",
                "client_context": {},
                "language": "de-DE",
                "constraints": {"max_length": 120, "required": True},
            },
            format="json",
        )
        self.assertEqual(res.status_code, 502)
        self.assertEqual(res.json().get("error_code"), "AI_PROVIDER_ERROR")
