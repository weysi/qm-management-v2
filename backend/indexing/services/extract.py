from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from zipfile import ZipFile

from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader


@dataclass
class ExtractedText:
    full_text: str
    slices: list[dict]


def _normalize(parts: list[str]) -> str:
    return "\n".join([part for part in parts if part and part.strip()]).strip()


def extract_pdf_text(path: Path) -> ExtractedText:
    reader = PdfReader(str(path))
    slices: list[dict] = []
    chunks: list[str] = []
    for index, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if not text:
            continue
        chunks.append(text)
        slices.append({"page": index, "text": text})
    return ExtractedText(full_text=_normalize(chunks), slices=slices)


def extract_docx_text(path: Path) -> ExtractedText:
    document = DocxDocument(str(path))
    paragraphs: list[str] = []
    slices: list[dict] = []

    for index, para in enumerate(document.paragraphs):
        text = para.text.strip()
        if not text:
            continue
        paragraphs.append(text)
        slices.append({"paragraph": index, "text": text})

    for table_index, table in enumerate(document.tables):
        for row_index, row in enumerate(table.rows):
            for col_index, cell in enumerate(row.cells):
                text = cell.text.strip()
                if not text:
                    continue
                paragraphs.append(text)
                slices.append(
                    {
                        "table": table_index,
                        "row": row_index,
                        "col": col_index,
                        "text": text,
                    }
                )

    return ExtractedText(full_text=_normalize(paragraphs), slices=slices)


def extract_pptx_text(path: Path) -> ExtractedText:
    presentation = Presentation(str(path))
    slides_text: list[str] = []
    slices: list[dict] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()
                if text:
                    parts.append(text)
                    slices.append({"slide": slide_index, "text": text})

        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                parts.append(notes_text)
                slices.append({"slide": slide_index, "notes": True, "text": notes_text})

        if parts:
            slides_text.append("\n".join(parts))

    return ExtractedText(full_text=_normalize(slides_text), slices=slices)


def extract_xlsx_text(path: Path) -> ExtractedText:
    wb = load_workbook(filename=str(path), data_only=True)
    lines: list[str] = []
    slices: list[dict] = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            row_text = " | ".join(
                [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
            )
            if not row_text:
                continue
            lines.append(f"[{sheet.title}] {row_text}")
            slices.append({"sheet": sheet.title, "text": row_text})
    return ExtractedText(full_text=_normalize(lines), slices=slices)


def extract_text_for_path(path: Path) -> ExtractedText:
    ext = path.suffix.lower()
    if ext == ".pdf":
        return extract_pdf_text(path)
    if ext in {".docx", ".doc"}:
        return extract_docx_text(path)
    if ext in {".pptx", ".ppt"}:
        return extract_pptx_text(path)
    if ext == ".xlsx":
        return extract_xlsx_text(path)

    # Fallback for plain text-like files.
    text = path.read_text(encoding="utf-8", errors="ignore")
    return ExtractedText(full_text=text.strip(), slices=[{"text": text.strip()}])


def extract_raw_ooxml_text(path: Path) -> str:
    """Best-effort raw OOXML text for placeholder scanning in split runs."""
    ext = path.suffix.lower()
    if ext not in {".docx", ".pptx", ".xlsx"}:
        return ""

    with ZipFile(path, "r") as archive:
        texts: list[str] = []
        for name in sorted(archive.namelist()):
            if not name.endswith(".xml"):
                continue
            if ext == ".docx" and not name.startswith("word/"):
                continue
            if ext == ".pptx" and not name.startswith("ppt/"):
                continue
            if ext == ".xlsx" and not name.startswith("xl/"):
                continue
            with archive.open(name) as handle:
                texts.append(handle.read().decode("utf-8", errors="ignore"))
        return "\n".join(texts)
